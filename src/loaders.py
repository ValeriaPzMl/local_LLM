import csv
import json
from pathlib import Path

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader


class DocumentLoader:
    TEXT_EXTENSIONS = {
        ".txt",
        ".md",
        ".py",
        ".js",
        ".ts",
        ".html",
        ".css",
        ".java",
        ".c",
        ".cpp",
        ".h",
        ".hpp",
        ".rs",
        ".go",
        ".sql",
        ".yaml",
        ".yml",
        ".toml",
        ".ini",
        ".log",
    }

    SUPPORTED_EXTENSIONS = (
        TEXT_EXTENSIONS
        | {
            ".pdf",
            ".docx",
            ".pptx",
            ".xlsx",
            ".csv",
            ".json",
        }
    )

    @classmethod
    def load(cls, path: str) -> str:
        file_path = Path(path)

        if not file_path.exists():
            raise FileNotFoundError(
                f"No existe el archivo: {file_path}"
            )

        extension = file_path.suffix.lower()

        if extension == ".pdf":
            return cls.load_pdf(file_path)

        if extension == ".docx":
            return cls.load_docx(file_path)

        if extension == ".pptx":
            return cls.load_pptx(file_path)

        if extension == ".xlsx":
            return cls.load_xlsx(file_path)

        if extension == ".csv":
            return cls.load_csv(file_path)

        if extension == ".json":
            return cls.load_json(file_path)

        if extension in cls.TEXT_EXTENSIONS:
            return cls.load_text(file_path)

        raise ValueError(
            f"Tipo de archivo no compatible: {extension}"
        )

    @staticmethod
    def load_pdf(path: Path) -> str:
        reader = PdfReader(path)
        pages = []

        for page_number, page in enumerate(
            reader.pages,
            start=1,
        ):
            text = page.extract_text()

            if text:
                pages.append(
                    f"[Página {page_number}]\n{text}"
                )

        return "\n\n".join(pages)

    @staticmethod
    def load_docx(path: Path) -> str:
        document = Document(path)
        blocks = []

        for paragraph in document.paragraphs:
            text = paragraph.text.strip()

            if text:
                blocks.append(text)

        for table_index, table in enumerate(
            document.tables,
            start=1,
        ):
            rows = []

            for row in table.rows:
                values = [
                    cell.text.strip()
                    for cell in row.cells
                ]

                rows.append(" | ".join(values))

            if rows:
                blocks.append(
                    f"[Tabla {table_index}]\n"
                    + "\n".join(rows)
                )

        return "\n\n".join(blocks)

    @staticmethod
    def load_pptx(path: Path) -> str:
        presentation = Presentation(path)
        slides = []

        for slide_number, slide in enumerate(
            presentation.slides,
            start=1,
        ):
            texts = []

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()

                    if text:
                        texts.append(text)

            if texts:
                slides.append(
                    f"[Diapositiva {slide_number}]\n"
                    + "\n".join(texts)
                )

        return "\n\n".join(slides)

    @staticmethod
    def load_xlsx(path: Path) -> str:
        workbook = load_workbook(
            path,
            read_only=True,
            data_only=True,
        )

        sheets = []

        for worksheet in workbook.worksheets:
            rows = []

            for row in worksheet.iter_rows(
                values_only=True
            ):
                values = [
                    "" if value is None else str(value)
                    for value in row
                ]

                if any(values):
                    rows.append(" | ".join(values))

            if rows:
                sheets.append(
                    f"[Hoja: {worksheet.title}]\n"
                    + "\n".join(rows)
                )

        workbook.close()

        return "\n\n".join(sheets)

    @staticmethod
    def load_csv(path: Path) -> str:
        rows = []

        with path.open(
            "r",
            encoding="utf-8",
            errors="replace",
            newline="",
        ) as file:
            reader = csv.reader(file)

            for row in reader:
                rows.append(" | ".join(row))

        return "\n".join(rows)

    @classmethod
    def load_json(cls, path: Path) -> str:
        with path.open(
            "r",
            encoding="utf-8",
            errors="replace",
        ) as file:
            data = json.load(file)

        lines = []
        cls.flatten_json(
            value=data,
            prefix="",
            output=lines,
        )

        return "\n".join(lines)


    @classmethod
    def flatten_json(
        cls,
        value,
        prefix: str,
        output: list[str],
    ) -> None:
        if isinstance(value, dict):
            for key, child_value in value.items():
                child_prefix = (
                    f"{prefix}.{key}"
                    if prefix
                    else str(key)
                )

                cls.flatten_json(
                    value=child_value,
                    prefix=child_prefix,
                    output=output,
                )

            return

        if isinstance(value, list):
            for index, child_value in enumerate(value):
                child_prefix = f"{prefix}[{index}]"

                cls.flatten_json(
                    value=child_value,
                    prefix=child_prefix,
                    output=output,
                )

            return

        output.append(
            f"{prefix}: {value}"
        )

    @staticmethod
    def load_text(path: Path) -> str:
        return path.read_text(
            encoding="utf-8",
            errors="replace",
        )